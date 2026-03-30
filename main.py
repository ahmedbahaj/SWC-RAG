"""
RAG over PowerPoint files in ./data using LangChain, Anthropic Haiku, ChromaDB,
and sentence-transformers/all-MiniLM-L6-v2 embeddings.
"""

from __future__ import annotations

import argparse
import os
from pathlib import Path
from typing import TYPE_CHECKING

from dotenv import load_dotenv
from langchain_anthropic import ChatAnthropic
from langchain_chroma import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation

if TYPE_CHECKING:
    from langchain_core.runnables import Runnable


DATA_DIR = Path(__file__).resolve().parent / "data"
CHROMA_DIR = Path(__file__).resolve().parent / "chroma_db"
COLLECTION_NAME = "pptx_rag"
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
DEFAULT_HAIKU = "claude-haiku-4-5"


def load_pptx_documents(data_dir: Path) -> list[Document]:
    paths = [
        p
        for p in sorted(data_dir.glob("**/*.pptx"))
        if not p.name.startswith("~$")
    ]
    if not paths:
        return []
    documents: list[Document] = []
    for pptx_path in paths:
        prs = Presentation(str(pptx_path))
        rel = pptx_path.relative_to(data_dir)
        for slide_idx, slide in enumerate(prs.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    t = shape.text.strip()
                    if t:
                        lines.append(t)
            content = "\n".join(lines)
            if content:
                documents.append(
                    Document(
                        page_content=content,
                        metadata={
                            "source": str(rel),
                            "slide": slide_idx,
                        },
                    )
                )
    return documents


def get_embeddings() -> HuggingFaceEmbeddings:
    return HuggingFaceEmbeddings(
        model_name=EMBEDDING_MODEL,
        model_kwargs={"device": "cpu"},
        encode_kwargs={"normalize_embeddings": True},
    )


def build_or_load_vectorstore(
    *,
    reindex: bool,
    data_dir: Path,
    persist_dir: Path,
) -> Chroma:
    embeddings = get_embeddings()
    if reindex and persist_dir.exists():
        import shutil

        shutil.rmtree(persist_dir)

    raw_docs = load_pptx_documents(data_dir)
    if not raw_docs:
        raise RuntimeError(
            f"No text extracted from PowerPoint files. Add .pptx files under {data_dir} "
            "and ensure slides contain extractable text."
        )

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
    )
    chunks = splitter.split_documents(raw_docs)

    if reindex or not persist_dir.exists():
        persist_dir.mkdir(parents=True, exist_ok=True)
        return Chroma.from_documents(
            documents=chunks,
            embedding=embeddings,
            persist_directory=str(persist_dir),
            collection_name=COLLECTION_NAME,
        )

    return Chroma(
        persist_directory=str(persist_dir),
        embedding_function=embeddings,
        collection_name=COLLECTION_NAME,
    )


def build_rag_chain(retriever, llm: ChatAnthropic):
    def format_docs(docs: list[Document]) -> str:
        parts = []
        for d in docs:
            src = d.metadata.get("source", "?")
            slide = d.metadata.get("slide", "?")
            parts.append(f"[{src} — slide {slide}]\n{d.page_content}")
        return "\n\n".join(parts)

    prompt = ChatPromptTemplate.from_messages(
        [
            (
                "system",
                "You answer questions using only the provided context from presentation slides. "
                "If the context is insufficient, say what is missing and answer only from what is given.",
            ),
            (
                "human",
                "Context:\n{context}\n\nQuestion: {question}",
            ),
        ]
    )

    return (
        RunnablePassthrough.assign(
            context=lambda x: format_docs(retriever.invoke(x["question"])),
        )
        | prompt
        | llm
        | StrOutputParser()
    )


def init_rag_chain(*, reindex: bool) -> "Runnable":
    load_dotenv()
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        raise RuntimeError("Set ANTHROPIC_API_KEY in your environment or .env file.")

    model_name = os.environ.get("ANTHROPIC_MODEL", DEFAULT_HAIKU)
    llm = ChatAnthropic(
        model=model_name,
        api_key=api_key,
        temperature=0,
    )

    vectorstore = build_or_load_vectorstore(
        reindex=reindex,
        data_dir=DATA_DIR,
        persist_dir=CHROMA_DIR,
    )
    retriever = vectorstore.as_retriever(search_kwargs={"k": 6})
    return build_rag_chain(retriever, llm)


def create_flask_app(chain: "Runnable"):
    from flask import Flask, jsonify, render_template, request

    app = Flask(__name__)

    @app.get("/")
    def index():
        return render_template("index.html")

    @app.post("/api/ask")
    def ask():
        body = request.get_json(silent=True) or {}
        question = (body.get("question") or "").strip()
        if not question:
            return jsonify({"error": "Missing or empty 'question' in JSON body."}), 400
        try:
            answer = chain.invoke({"question": question})
            return jsonify({"answer": answer})
        except Exception as e:
            return jsonify({"error": str(e)}), 500

    return app


def main() -> None:
    parser = argparse.ArgumentParser(description="RAG over PowerPoint files in ./data")
    parser.add_argument(
        "--reindex",
        action="store_true",
        help="Rebuild the Chroma index from ./data (deletes existing persist dir).",
    )
    parser.add_argument(
        "-q",
        "--question",
        type=str,
        default=None,
        help="Ask a single question and exit (otherwise interactive REPL).",
    )
    parser.add_argument(
        "--serve",
        action="store_true",
        help="Run a small web UI (HTML) and JSON API instead of the terminal REPL.",
    )
    parser.add_argument(
        "--host",
        type=str,
        default="127.0.0.1",
        help="Bind address for --serve (default: 127.0.0.1).",
    )
    parser.add_argument(
        "--port",
        type=int,
        default=5000,
        help="Port for --serve (default: 5000).",
    )
    args = parser.parse_args()

    try:
        chain = init_rag_chain(reindex=args.reindex)
    except RuntimeError as e:
        raise SystemExit(str(e)) from e

    if args.serve:
        app = create_flask_app(chain)
        print(f"Open http://{args.host}:{args.port}/ in your browser (Ctrl+C to stop).")
        app.run(host=args.host, port=args.port, debug=False)
        return

    if args.question:
        out = chain.invoke({"question": args.question})
        print(out)
        return

    print("RAG ready. Type a question (empty line to quit).")
    while True:
        try:
            q = input("\nQuestion: ").strip()
        except (EOFError, KeyboardInterrupt):
            print()
            break
        if not q:
            break
        print(chain.invoke({"question": q}))


if __name__ == "__main__":
    main()
